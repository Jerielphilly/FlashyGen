import re # Import the re module for regular expressions
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
import fitz # PyMuPDF for PDF processing
from pptx import Presentation # For PPTX processing
import ollama # Ollama Python client for LLM interaction
import io # For handling file-like objects
import json

# Initialize FastAPI app
app = FastAPI()

# Configure CORS middleware
# This allows your frontend (e.g., running on http://localhost:5173) to make requests to this backend.
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173"], # Specify your frontend's origin
    allow_credentials=True,
    allow_methods=["*"], # Allow all HTTP methods (GET, POST, etc.)
    allow_headers=["*"], # Allow all headers
)

# Function to extract text from a PDF file
async def extract_text_from_pdf(file: UploadFile) -> str:
    """
    Extracts text content from an uploaded PDF file.
    """
    try:
        file_bytes = await file.read() # Read the file content as bytes
        # Open the PDF from bytes in memory
        with fitz.open(stream=file_bytes, filetype="pdf") as doc:
            text = ""
            for page in doc:
                text += page.get_text() # Extract text from each page
        return text
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process PDF: {e}")

# Function to extract text from a PPTX file
def extract_text_from_pptx(file_obj) -> str:
    """
    Extracts text content from an uploaded PPTX file.
    Expects a file-like object (e.g., file.file from UploadFile).
    """
    text = ""
    try:
        prs = Presentation(file_obj) # Load the presentation
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): # Check if the shape contains text
                    text += shape.text + "\n" # Append text and a newline
        return text
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process PPTX: {e}")

# Function to generate flashcards using Ollama (Mistral model)
def generate_flashcards_with_mistral(text: str) -> list:
    """
    Generates flashcards (question-answer pairs) from the given text
    using the 'mistral' model via Ollama.
    """
    # Define the prompt for the LLM to generate flashcards in JSON format
    # Emphasize strict JSON output without extra text or markdown wrappers
    prompt = f"""
You are an expert AI that generates flashcards for efficient learning.

Given the following content, your job is to create a **comprehensive set of high-quality question-answer flashcards**. Follow these rules:

1. Break down all **topics, subtopics, and key points** into distinct flashcards.
2. Ensure each flashcard captures **only one concept** or fact.
3. Use **clear, specific questions** that test understanding, not just memory.
4. Provide **complete, concise answers** — avoid vague responses.
5. Include **important definitions, examples, steps, and reasoning** wherever relevant.
6. Avoid duplicating questions, and skip filler content.
7. Do **not limit yourself to a fixed number of flashcards** — generate as many as needed to cover the material completely.
8. **CRITICAL: Output ONLY the JSON array. Do NOT include any conversational text, explanations, or markdown code block delimiters (like ```json or ```) before or after the JSON array.**

### Input Content:
\"\"\"
{text[:4000]}
\"\"\"

### Output Format (Strict JSON array, nothing else):
[
    {{
    "question": "What is ...?",
    "answer": "..."
    }},
    ...
]
"""

    try:
        # Call Ollama's chat API
        response = ollama.chat(model="mistral", messages=[{"role": "user", "content": prompt}])
        # Extract the content from the response
        generated_content = response['message']['content']

        # --- NEW: Extract only the JSON part using regex ---
        # This regex looks for a JSON array (starts with '[', ends with ']')
        # and captures everything in between. It handles cases where there might be
        # leading/trailing whitespace or markdown code blocks.
        json_match = re.search(r'\[.*\]', generated_content, re.DOTALL)
        if json_match:
            json_string = json_match.group(0)
        else:
            # Fallback if regex doesn't find a perfect match, try to clean
            # This might still fail if the structure is completely off
            json_string = generated_content.strip().replace('```json', '').replace('```', '').strip()
            print(f"Warning: Regex failed to find JSON. Attempting fallback cleaning. Cleaned content: {json_string[:200]}...")


        # Attempt to parse the content as a Python list (from JSON string)
        flashcards = json.loads(json_string)
        return flashcards
    except json.JSONDecodeError as e:
        print(f"Error decoding JSON from Ollama response: {e}")
        print(f"Ollama's raw response content: {generated_content}") # Still useful for debugging what LLM sent
        print(f"Attempted JSON string for parsing: {json_string}") # Show what we tried to parse
        raise HTTPException(status_code=500, detail="Ollama did not return valid JSON for flashcards after extraction.")
    except Exception as e:
        print(f"Error generating flashcards with Ollama: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate flashcards: {e}")

# FastAPI endpoint for file upload and quiz generation
@app.post("/upload/")
async def upload_file(file: UploadFile = File(...)):
    """
    Handles file uploads (PDF, PPTX, TXT), extracts text, and generates flashcards.
    """
    ext = file.filename.split(".")[-1].lower()
    text = ""

    # Extract text based on file type
    if ext == "pdf":
        text = await extract_text_from_pdf(file)
    elif ext == "pptx":
        # For python-pptx, we need a file-like object, not raw bytes directly from UploadFile
        # So, we read the bytes and then create an in-memory BytesIO object
        file_bytes = await file.read()
        text = extract_text_from_pptx(io.BytesIO(file_bytes))
    elif ext == "txt":
        text = (await file.read()).decode("utf-8")
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type. Only PDF, PPTX, or TXT are allowed.")

    if not text.strip(): # Check if extracted text is empty after stripping whitespace
        raise HTTPException(status_code=400, detail="Could not extract content from the file or file is empty.")

    # Generate flashcards using the extracted text
    flashcards = generate_flashcards_with_mistral(text)

    # Ensure only 'question' and 'answer' fields are returned and filter out malformed entries
    filtered_flashcards = [
        {"question": fc.get("question", ""), "answer": fc.get("answer", "")}
        for fc in flashcards if isinstance(fc, dict) and "question" in fc and "answer" in fc
    ]

    if not filtered_flashcards:
        raise HTTPException(status_code=500, detail="No valid flashcards could be generated from the content.")

    return {"quiz": filtered_flashcards}