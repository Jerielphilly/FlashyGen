import fitz
from pptx import Presentation

def extract_text_from_pdf(file):
    doc = fitz.open(stream=file.file.read(), filetype="pdf")
    return "\n".join([page.get_text() for page in doc])

def extract_text_from_pptx(file):
    prs = Presentation(file.file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_txt(file):
    return file.file.read().decode("utf-8")
